#!/usr/bin/env python3
"""Build premium 21-slide presentation."""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "src"))

from mee344_g05.config import FIGURES_DIR, REPORTS_DIR, SLIDES_DIR

# Colors
C_NAVY = RGBColor(0x0B, 0x1F, 0x3A)
C_NAVY2 = RGBColor(0x13, 0x2D, 0x50)
C_TEAL = RGBColor(0x00, 0xC9, 0xA7)
C_SKY = RGBColor(0x38, 0xBD, 0xF8)
C_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_OFF = RGBColor(0xF8, 0xFA, 0xFC)
C_SLATE = RGBColor(0x47, 0x55, 0x69)
C_CARD = RGBColor(0xEE, 0xF2, 0xF7)

SW = Inches(13.333)
SH = Inches(7.5)
HEADER = Inches(0.9)
M = Inches(0.42)
CT = HEADER + Inches(0.15)

FIG = FIGURES_DIR / "presentation"
BG = ROOT / "assets" / "backgrounds"

MEMBERS = [
    "Berhan Kaya", "Ahmet Bayram Topçu", "Yusuf Duman",
    "Furkan Efe Demirbel", "Tekgül Eroğlu",
]

AGENDA = [
    ("01", "Problem & Veri", "Berhan", "2 dk"),
    ("02", "EDA & Pipeline", "Ahmet", "3 dk"),
    ("03", "Features & Modeller", "Yusuf", "2 dk"),
    ("04", "Sonuçlar", "Furkan", "3 dk"),
    ("05", "Yorum & Demo", "Tekgül", "3 dk"),
]


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(sl, l, t, w, h, fill, line=None):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s


def rnd(sl, l, t, w, h, fill):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def txt(sl, l, t, w, h, text, size=14, bold=False, color=C_SLATE, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(l, t, w, h)
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    tb.text_frame.word_wrap = True
    return tb


def bg_image(sl, path: Path):
    if path.exists():
        sl.shapes.add_picture(str(path), 0, 0, width=SW, height=SH)
    else:
        rect(sl, 0, 0, SW, SH, C_NAVY2)


def header(sl, title, sub=""):
    rect(sl, 0, 0, SW, HEADER, C_NAVY)
    rect(sl, 0, HEADER - Inches(0.05), SW, Inches(0.05), C_TEAL)
    txt(sl, M, Inches(0.16), Inches(10), Inches(0.45), title, 26, True, C_WHITE)
    if sub:
        txt(sl, M, Inches(0.58), Inches(9), Inches(0.28), sub, 11, False, C_SKY)
    txt(sl, SW - Inches(2), Inches(0.25), Inches(1.6), Inches(0.3), "G05 · MEE344", 10, False, C_TEAL, PP_ALIGN.RIGHT)


def footer(sl, note="EPİAŞ STLF · Decision Tree + AdaBoost"):
    rect(sl, 0, SH - Inches(0.35), SW, Inches(0.35), C_CARD)
    txt(sl, M, SH - Inches(0.3), Inches(11), Inches(0.22), note, 9, False, C_SLATE)


def bullets(sl, l, t, w, h, items, size=12):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = C_SLATE
        p.space_after = Pt(6)


def pic(sl, path, l, t, w):
    if path and Path(path).exists():
        sl.shapes.add_picture(str(path), l, t, width=w)


def kpi(sl, l, t, w, h, label, value, sub="", accent=C_TEAL):
    rnd(sl, l, t, w, h, C_WHITE)
    rect(sl, l, t, w, Inches(0.07), accent)
    txt(sl, l + Inches(0.1), t + Inches(0.12), w, Inches(0.28), label, 9, False, C_SLATE)
    txt(sl, l + Inches(0.1), t + Inches(0.38), w, Inches(0.42), value, 20, True, C_NAVY)
    if sub:
        txt(sl, l + Inches(0.1), t + Inches(0.78), w, Inches(0.22), sub, 8, False, C_SLATE)


def section(prs, num, title, desc, bg_name):
    s = blank(prs)
    bg_image(s, BG / bg_name)
    txt(s, M, Inches(2.0), Inches(2.5), Inches(1.2), num, 96, True, C_TEAL)
    txt(s, M + Inches(2.2), Inches(2.4), Inches(9), Inches(0.7), title, 34, True, C_WHITE)
    txt(s, M + Inches(2.2), Inches(3.2), Inches(9), Inches(0.5), desc, 16, False, C_SKY)


def slide_cover(prs, metrics):
    s = blank(prs)
    bg_image(s, BG / "bg_cover.png")
    rnd(s, M, Inches(0.35), Inches(2.2), Inches(0.4), C_TEAL)
    txt(s, M + Inches(0.15), Inches(0.42), Inches(2), Inches(0.3), "MEE344 · GRUP G05", 10, True, C_NAVY)
    txt(s, M, Inches(1.0), Inches(8.5), Inches(1.1), "Türkiye Saatlik Elektrik\nTüketim Tahmini", 38, True, C_WHITE)
    txt(s, M, Inches(2.35), Inches(8), Inches(0.45), "Decision Tree + AdaBoost · Short-Term Load Forecasting", 17, False, C_SKY)
    y = Inches(3.2)
    for i, name in enumerate(MEMBERS):
        col, row = i % 3, i // 3
        rnd(s, M + col * Inches(2.15), y + row * Inches(0.5), Inches(2), Inches(0.38), C_NAVY2)
        txt(s, M + col * Inches(2.15) + Inches(0.08), y + row * Inches(0.5) + Inches(0.06), Inches(1.85), Inches(0.28), name, 10, False, C_WHITE)
    ada = metrics["test"]["adaboost"]
    kpi(s, Inches(8.8), Inches(4.5), Inches(3.8), Inches(1.2), "En iyi model (test)",
        f"R² {ada['r2']:.3f}", f"RMSE {ada['rmse']:.0f} MWh · MAPE {ada['mape']:.1f}%")
    txt(s, M, SH - Inches(0.55), Inches(6), Inches(0.3), "22 Mayıs 2026 · EPİAŞ 31.07–29.10.2025", 10, False, C_SKY)


def slide_agenda(prs):
    s = blank(prs)
    header(s, "Sunum Akışı", "10–15 dakika · 5 konuşmacı")
    footer(s)
    rect(s, M, CT, SW - 2 * M, Inches(0.08), C_TEAL)
    for i, (num, title, who, dur) in enumerate(AGENDA):
        top = CT + Inches(0.2) + i * Inches(1.05)
        rnd(s, M, top, SW - 2 * M, Inches(0.88), C_OFF if i % 2 else C_WHITE)
        txt(s, M + Inches(0.15), top + Inches(0.2), Inches(0.5), Inches(0.5), num, 18, True, C_TEAL)
        txt(s, M + Inches(0.7), top + Inches(0.15), Inches(5), Inches(0.4), title, 15, True, C_NAVY)
        txt(s, M + Inches(0.7), top + Inches(0.48), Inches(3), Inches(0.3), who, 11, False, C_SLATE)
        txt(s, SW - Inches(1.8), top + Inches(0.28), Inches(1.2), Inches(0.35), dur, 11, False, C_SKY, PP_ALIGN.RIGHT)


def slide_model_importance(prs, test: dict) -> None:
    """Decision Tree + AdaBoost — sol metin, sağda importance grafikleri."""
    s = blank(prs)
    dt, ada = test["decision_tree"], test["adaboost"]
    imp_pct = (1 - ada["rmse"] / dt["rmse"]) * 100
    header(
        s,
        "Decision Tree & AdaBoost",
        f"Test RMSE DT {dt['rmse']:.0f} MWh · AdaBoost {ada['rmse']:.0f} MWh",
    )
    footer(s)
    rnd(s, M, CT, Inches(5.55), Inches(5.35), C_OFF)
    bullets(
        s,
        M + Inches(0.15),
        CT + Inches(0.12),
        Inches(5.25),
        Inches(5.1),
        [
            "Decision Tree — tek ağaç, yorumlanabilir baseline",
            "Neredeyse tüm karar: lag_168h (geçen hafta aynı saat)",
            "lag_1h ikincil; üretim / saat özellikleri düşük MDI",
            "",
            "AdaBoost — 200 zayıf öğrenici → güçlü ensemble",
            "lag_168h + lag_1h birlikte (haftalık + saatlik döngü)",
            f"RMSE: {dt['rmse']:.0f} → {ada['rmse']:.0f} MWh (~%{imp_pct:.0f} iyileşme)",
            f"R²: {dt['r2']:.3f} → {ada['r2']:.3f} · MAPE {ada['mape']:.2f}%",
        ],
        11,
    )
    panel = FIG / "dt_ada_importance_panel.png"
    if panel.exists():
        pic(s, panel, Inches(5.95), CT - Inches(0.05), Inches(7.0))
    else:
        pic(s, FIG / "dt_importance.png", Inches(6.1), CT, Inches(6.9))
        pic(s, FIG / "ada_importance.png", Inches(6.1), CT + Inches(2.75), Inches(6.9))
        txt(
            s,
            Inches(6.1),
            CT + Inches(2.48),
            Inches(6.9),
            Inches(0.25),
            "Üst: Decision Tree · Alt: AdaBoost — özellik önemi (MDI)",
            9,
            False,
            C_SLATE,
            PP_ALIGN.CENTER,
        )


def content(prs, title, sub, items, img_name, img_w=Inches(6)):
    s = blank(prs)
    header(s, title, sub)
    footer(s)
    bullets(s, M, CT + Inches(0.05), Inches(5.6), Inches(5.5), items)
    pic(s, FIG / img_name, SW - img_w - M, CT, img_w)


def slide_web_demo(prs):
    s = blank(prs)
    header(s, "Canlı Tahmin — Web Demo", "FastAPI · Tarayıcıda çalışan ML")
    footer(s)
    shots = [
        ("web_overview.png", "Ana panel"),
        ("web_backtest.png", "Backtest"),
        ("web_forecast.png", "İleri tahmin"),
    ]
    w = Inches(3.9)
    for i, (fname, cap) in enumerate(shots):
        left = M + i * (w + Inches(0.2))
        p = ROOT / "assets" / "screenshots" / fname
        if p.exists():
            pic(s, p, left, CT + Inches(0.1), w)
        else:
            rnd(s, left, CT + Inches(0.5), w, Inches(2.5), C_CARD)
            txt(s, left, CT + Inches(1.4), w, Inches(0.4), cap, 12, False, C_SLATE, PP_ALIGN.CENTER)
        txt(s, left, CT + Inches(3.0), w, Inches(0.3), cap, 10, False, C_SLATE, PP_ALIGN.CENTER)
    rnd(s, M, Inches(5.8), Inches(4), Inches(0.45), C_TEAL)
    txt(s, M + Inches(0.12), Inches(5.88), Inches(3.8), Inches(0.35), "Hiçbir grupta yok — canlı API + grafik", 10, True, C_NAVY)


def slide_thanks(prs):
    s = blank(prs)
    bg_image(s, BG / "bg_thanks.png")
    txt(s, M, Inches(2.6), SW - 2 * M, Inches(0.9), "Teşekkürler", 48, True, C_WHITE, PP_ALIGN.CENTER)
    txt(s, M, Inches(3.6), SW - 2 * M, Inches(0.5), "Sorularınız?", 28, False, C_TEAL, PP_ALIGN.CENTER)
    txt(s, M, Inches(4.5), SW - 2 * M, Inches(1), "\n".join(MEMBERS), 12, False, C_SKY, PP_ALIGN.CENTER)
    txt(s, M, Inches(6.2), SW - 2 * M, Inches(0.35), "github.com · mee344-g05-stlf", 11, False, C_WHITE, PP_ALIGN.CENTER)


def metrics_table_slide(prs, metrics):
    s = blank(prs)
    header(s, "Detaylı Metrik Tablosu", "Hold-out test · 336 saat")
    footer(s)
    test = metrics["test"]
    rows = [
        ("Model", "RMSE", "MAE", "R²", "MAPE"),
        ("Naive 1h", f"{test['baseline_last_hour']['rmse']:.0f}", f"{test['baseline_last_hour']['mae']:.0f}",
         f"{test['baseline_last_hour']['r2']:.3f}", f"{test['baseline_last_hour']['mape']:.2f}%"),
        ("Seasonal 168h", f"{test['baseline_seasonal_168h']['rmse']:.0f}", f"{test['baseline_seasonal_168h']['mae']:.0f}",
         f"{test['baseline_seasonal_168h']['r2']:.3f}", f"{test['baseline_seasonal_168h']['mape']:.2f}%"),
        ("Decision Tree", f"{test['decision_tree']['rmse']:.0f}", f"{test['decision_tree']['mae']:.0f}",
         f"{test['decision_tree']['r2']:.3f}", f"{test['decision_tree']['mape']:.2f}%"),
        ("AdaBoost", f"{test['adaboost']['rmse']:.0f}", f"{test['adaboost']['mae']:.0f}",
         f"{test['adaboost']['r2']:.3f}", f"{test['adaboost']['mape']:.2f}%"),
    ]
    tw = SW - 2 * M
    th = Inches(2.2)
    cw = tw / 5
    top = CT + Inches(0.1)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            fill = C_NAVY if r == 0 else (RGBColor(0xE0, 0xFB, 0xF4) if r == 4 else (C_OFF if r % 2 else C_WHITE))
            col = C_WHITE if r == 0 else C_SLATE
            rnd(s, M + c * cw + Inches(0.02), top + r * (th / 5) + Inches(0.02),
                cw - Inches(0.04), th / 5 - Inches(0.04), fill)
            txt(s, M + c * cw + Inches(0.05), top + r * (th / 5) + Inches(0.04),
                cw - Inches(0.08), th / 5 - Inches(0.06), val, 10 if r else 11,
                r == 0 or r == 4, col, PP_ALIGN.CENTER)
    pic(s, FIG / "actual_vs_pred_both.png", M, CT + Inches(2.5), SW - 2 * M)


def kpi_dashboard(prs, metrics):
    s = blank(prs)
    header(s, "Model Karşılaştırması", "Test seti performansı")
    footer(s)
    test = metrics["test"]
    cards = [
        ("AdaBoost RMSE", f"{test['adaboost']['rmse']:.0f}", "MWh", C_TEAL),
        ("AdaBoost R²", f"{test['adaboost']['r2']:.3f}", "test", C_TEAL),
        ("DT RMSE", f"{test['decision_tree']['rmse']:.0f}", "MWh", C_SKY),
        ("MAPE", f"{test['adaboost']['mape']:.2f}%", "AdaBoost", C_AMBER),
    ]
    cw = Inches(2.85)
    for i, (lbl, val, sub, col) in enumerate(cards):
        kpi(s, M + i * (cw + Inches(0.18)), CT, cw, Inches(1.1), lbl, val, sub, col)
    pic(s, FIG / "metrics_rmse_comparison.png", M, CT + Inches(1.3), SW - 2 * M)


def main():
    metrics = json.loads((REPORTS_DIR / "metrics.json").read_text(encoding="utf-8"))
    test = metrics["test"]

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    slide_cover(prs, metrics)
    slide_agenda(prs)

    section(prs, "01", "Problem Tanımı", "STLF · Regression · EPİAŞ verisi", "bg_section_01.png")
    content(prs, "Problem & Motivasyon", "Saatlik elektrik tüketimi (MWh)",
            [
                "STLF: kısa vadeli talep tahmini — üretim planlama ve dengeleme",
                "Regression: hedef consumption_mwh, metrik RMSE / MAE / R²",
                "Üretim kaynakları exogenous feature (leakage yok)",
                "G05 allocation: Decision Tree + AdaBoost",
            ], "consumption_series.png")

    section(prs, "02", "Veri & Keşif", "EDA · 2169 saat", "bg_section_02.png")
    content(prs, "Veri Seti", "31.07.2025 – 29.10.2025",
            [
                "Kaynak: EPİAŞ Şeffaflık Platformu",
                "Üretim: 16 kaynak + toplam | Tüketim: MWh",
                "Train 1665 h | Test 336 h (son 14 gün)",
                "26 feature — cyclic, lag, rolling, exogenous",
            ], "consumption_series.png")
    content(prs, "EDA — Tüketim Örüntüleri", "Günlük ve haftalık siklus",
            [
                "Akşam piki 18–21 saatleri",
                "Hafta sonu tüketim farkı",
                "Tatil günlerinde düşüş (30 Ağustos, 29 Ekim)",
                "Lag 24h ve 168h güçlü otokorelasyon",
            ], "weekly_heatmap.png")
    content(prs, "EDA — Üretim Karışımı", "Termik vs yenilenebilir",
            ["Doğal gaz ve kömür baskın", "Güneş gün ortası artışı", "Rüzgar değişken"],
             "production_mix.png")

    section(prs, "03", "Pipeline & Modeller", "Feature engineering · CV", "bg_section_03.png")
    s = blank(prs)
    header(s, "Veri Pipeline", "Uçtan uca reproducible")
    footer(s)
    pic(s, FIG / "pipeline_horizontal.png", M, CT + Inches(0.3), SW - 2 * M)
    bullets(s, M, Inches(5.5), Inches(10), Inches(1.2),
            ["Locale parse (TR ondalık)", "TimeSeriesSplit(5) + GridSearchCV", "random_state=42"], 11)

    content(prs, "Feature Engineering", "Leakage-safe",
            [
                "Cyclic: hour/dow/doy sin-cos",
                "Lag: 1h, 3h, 24h, 168h",
                "Rolling: mean/std (shift 1)",
                "Exogenous: wind, solar, hydro, thermal",
            ], "cyclic_demo.png")
    s = blank(prs)
    header(s, "Modelleme", "Decision Tree + AdaBoost")
    footer(s)
    pic(s, FIG / "weak_to_strong.png", M, CT + Inches(0.2), Inches(5.5))
    bullets(s, Inches(6.2), CT + Inches(0.3), Inches(6.5), Inches(4),
            ["DT: yorumlanabilir non-linear baseline",
             "AdaBoost: 200 stump, square loss",
             f"CV RMSE DT {metrics['cv']['decision_tree']:.0f} | Ada {metrics['cv']['adaboost']:.0f}"],
            12)

    section(prs, "04", "Sonuçlar", "Test performansı", "bg_section_04.png")
    kpi_dashboard(prs, metrics)
    metrics_table_slide(prs, metrics)
    slide_model_importance(prs, test)
    s = blank(prs)
    header(s, "Tahmin Kalitesi", "14 günlük hold-out")
    footer(s)
    pic(s, FIG / "actual_vs_pred_both.png", M, CT, SW - 2 * M)

    section(prs, "05", "Yorum & Demo", "Importance · Web", "bg_section_05.png")
    content(prs, "Yorumlama & Hata", "Feature importance · residual",
            ["Top: lag_24h, lag_168h, hour cycle", "Akşam saatlerinde artık hata",
             "Permutation importance doğrulandı"],
             "residual_hour.png" if (FIG / "residual_hour.png").exists() else "ada_importance.png")
    slide_web_demo(prs)
    content(prs, "Sınırlılıklar & Gelecek", "",
            ["Sınır: hava verisi yok, 3 ay pencere", "Sınır: lag → gerçek zaman gecikme",
             "Gelecek: hava API, XGBoost, ensemble", "Gelecek: EPİAŞ canlı feed"],
             "metrics_rmse_comparison.png")
    slide_thanks(prs)

    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    out = SLIDES_DIR / "G05_STLF_DT_AdaBoost.pptx"
    try:
        prs.save(out)
    except PermissionError:
        out = SLIDES_DIR / "G05_STLF_DT_AdaBoost_NEW.pptx"
        prs.save(out)
        print("Ana dosya açık — yeni dosyaya kaydedildi.")
    print(f"Saved {len(prs.slides)} slides -> {out}")


if __name__ == "__main__":
    main()
